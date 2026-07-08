"""EA Slide Builder self-test.

Run from the project folder (venv active):

    python -m tests.run_all

Uses only the app's own dependencies (no pytest). Exercises the parsers, the
date/phase math, the slide/deck builders, the profiles round-trip, and finally
the full Streamlit script via AppTest — including the Generate path.
Exits non-zero on any failure.
"""

from __future__ import annotations

import sys
import tempfile
from datetime import date
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

FAILURES: list[str] = []


def check(name: str, cond: bool, detail: str = "") -> None:
    status = "ok  " if cond else "FAIL"
    print(f"[{status}] {name}" + (f" — {detail}" if detail and not cond else ""))
    if not cond:
        FAILURES.append(name)


# --------------------------------------------------------------------------- #
# Parsers & computations
# --------------------------------------------------------------------------- #
def test_parsers():
    import data_processor as dp

    # Long (Tableau) machine format, with a trailing partial quarter to trim.
    long_text = (
        "Year\tQuarter\tMonth\tMachine Type\tDistinct count of machine_id\n"
        "2025\tQ1\tJanuary\tExisting\t100\n"
        "2025\tQ1\tFebruary\tExisting\t110\n"
        "2025\tQ1\tMarch\tExisting\t120\n"
        "2025\tQ1\tJanuary\tNew\t10\n"
        "2025\tQ1\tFebruary\tNew\t20\n"
        "2025\tQ1\tMarch\tNew\t30\n"
        "2025\tQ2\tApril\tExisting\t130\n"
        "2025\tQ2\tMay\tExisting\t140\n"
        "2025\tQ2\tJune\tExisting\t150\n"
        "2026\tQ1\tJanuary\tExisting\t40\n"
    )
    df = dp.parse_machine_count(long_text)
    check("machine long-format rows", len(df) == 2, f"got {len(df)}")
    check("machine quarterly average",
          int(df.iloc[0]["total"]) == 130,  # avg(100,110,120)+avg(10,20,30)
          f"got {df.iloc[0]['total']}")
    check("partial quarter trimmed", "2026" not in " ".join(df["period"]))

    # Wide fallback.
    wide = dp.parse_machine_count("P\tNew\tExisting\nQ1\t10\t100\nQ2\t20\t120")
    check("machine wide-format rows", len(wide) == 2)
    stats = dp.compute_machine_stats(wide)
    check("stats max/min", stats["max_total"] == 140 and stats["min_total"] == 110)

    # Geo locations default to avoiding product double-counting + state abbreviation.
    geo = ("ip_country\tip_region\tip_city\tMeasure Names\tproduct_name\tMeasure Values\n"
           "United States\tConnecticut\tBristol\tDistinct count\tLabVIEW\t1\n"
           "United States\tTexas\tAustin\tDistinct count\tLabVIEW\t50\n"
           "United States\tTexas\tAustin\tDistinct count\tDIAdem\t30\n")
    ldf = dp.parse_locations(geo)
    austin = ldf[ldf["city"] == "Austin"].iloc[0]
    check("geo avoids product double-counting",
          int(austin["count"]) == 50, f"got {austin['count']}")
    summed = dp.parse_locations(geo, avoid_product_double_count=False)
    summed_austin = summed[summed["city"] == "Austin"].iloc[0]
    check("geo can sum product rows",
          int(summed_austin["count"]) == 80, f"got {summed_austin['count']}")
    check("state abbreviation",
          ldf[ldf["city"] == "Bristol"].iloc[0]["state"] == "CT")
    geo_with_subtotals = (
        "ip_country\tip_region\tip_city\tMeasure Names\tproduct_name\tMeasure Values\n"
        "United States\tConnecticut\t\tDistinct count of machine_id\t*\t2,448.00\n"
        "United States\t\t\tDistinct count of machine_id\t*\t1,287.00\n"
        "United States\tCalifornia\tOntario\tDistinct count of machine_id\t*\t983\n"
        "United States\tIowa\tMarion\tDistinct count of machine_id\t*\t177\n"
        "United States\tArizona\tTucson\tDistinct count of machine_id\t*\t37\n"
        "United States\tCalifornia\tLos Angeles\tDistinct count of machine_id\t*\t25\n"
        "United States\tIowa\tCedar Rapids\tDistinct count of machine_id\t*\t15\n"
        "United States\tTexas\tQ1 2025\tDistinct count of machine_id\t*\t999\n"
    )
    top_sites = dp.top_locations(dp.parse_locations(geo_with_subtotals), 5)
    check("geo ignores subtotal/non-site rows",
          top_sites[0]["city"] == "Ontario" and len(top_sites) == 5,
          f"got {top_sites}")

    # Versions: top version per product.
    vdf = dp.parse_usage_versions(
        "product_name\tproduct_version\tcount\n"
        "LabVIEW\t2021\t399\nLabVIEW\t2023\t342\nMAX\t20.0.0\t479\n")
    top = dp.top_versions(vdf)
    check("top version", top[0]["version"] == "2021" and top[0]["users"] == 399)
    check("top software total and share",
          top[0]["product"] == "LabVIEW"
          and top[0]["product_total"] == 741
          and top[0]["pct"] == 54,
          f"got {top[0]}")

    # Date/phase math.
    start = dp.parse_date("02-JAN-2026")
    check("date parse", start == date(2026, 1, 2))
    end = dp.compute_end_date(start, dp.parse_term_years("3 years"))
    check("inclusive end date", end == date(2029, 1, 1), f"got {end}")
    ph = dp.compute_phase(start, 3.0, today=date(2027, 9, 1))
    check("phase second half", ph["phase"] == "Second Half"
          and ph["hint"] == "Year 2 of 3")
    check("phase not started",
          dp.compute_phase(start, 3.0, today=date(2025, 1, 1))["phase"]
          == "Not started")
    check("phase expired",
          dp.compute_phase(start, 3.0, today=date(2030, 1, 1))["phase"]
          == "Expired")


def test_screenshot_parsers():
    import screenshot_reader as sr

    contract = ("EA/EP Service ID 1   EA-15725\n"
                "Company 2   RTX Corporation\n"
                "Effective Date   02-JAN-2026\n"
                "EP Term   3 years\n"
                "EA FLEX Credits   20,000\n"
                "Support Level   Enterprise Support\n"
                "Debug licenses included   Yes\n")
    p = sr.parse_contract_details(contract)
    check("contract service id", p["service_id"] == "EA-15725")
    check("contract customer", p["customer"] == "RTX Corporation")
    check("contract support", p["support_level"] == "Enterprise Support")
    check("contract debug", p["debug_licenses"] == "Yes")

    finite = ("Finite Quantity NI Software Licenses\n"
              "142 Named User or Computer Based\tLabVIEW Full\n"
              "23 Concurrent\tTestStand\n")
    rows = sr.parse_finite_licenses(finite)
    check("finite rows", len(rows) == 2)
    check("finite split", rows[0]["count"] == 142
          and rows[0]["license_name"] == "LabVIEW Full")

    bundles = sr.parse_unlimited_bundles(
        "Unlimited Quantity NI Software Licenses\n"
        "Unlimited Named-User\tEP Bundle Title: EA Platform Bundle\n")
    check("bundle name", bundles == ["EA Platform Bundle"])

    conf = sr.field_confidence(
        [{"text": "EA-15725", "conf": 91.0}], "EA-15725")
    check("field confidence", conf == 91.0, f"got {conf}")


# --------------------------------------------------------------------------- #
# Slide / deck / preview / profiles
# --------------------------------------------------------------------------- #
def _sample_data():
    import data_processor as dp
    mdf = dp.parse_machine_count(
        "P\tNew\tExisting\n" + "\n".join(
            f"Q{i % 4 + 1} {2023 + i // 4}\t{50 + i}\t{300 + i * 20}"
            for i in range(8)))
    return {
        "service_id": "EA-15725", "customer": "Blue Origin",
        "updated_date": "16-Jun-2026", "ea_end_date": "01-JAN-2029",
        "ep_term": "3 years", "contract_scope": "All NI software",
        "phase": "Second Half",
        "bundles": ["EA Platform Bundle"],
        "finite_licenses": [{"count": 142, "license_type": "Named User",
                             "license_name": "LabVIEW Full"}],
        "machine": {"df": mdf, "stats": dp.compute_machine_stats(mdf)},
        "locations_top5": [{"state": "TX", "city": "Austin",
                            "location": "Austin, TX", "count": 820}],
        "versions_top5": [{"product": "LabVIEW", "version": "2021",
                           "users": 399, "product_total": 2145, "pct": 19}],
        "credits": {"purchased": "20,000", "used": "7,500", "pct_used": 38},
        "support": {"tier": "Enterprise Support", "scope": "All Users"},
    }


def test_slide_builder():
    from pptx import Presentation
    from insights import generate_insights
    from slide_builder import build_deck, build_slide

    data = _sample_data()
    ins = generate_insights(data)
    check("insights generated", len(ins) >= 3, f"got {len(ins)}")

    buf = build_slide(data, insights=ins)
    prs = Presentation(buf)
    check("slide + insights = 2 slides", len(list(prs.slides)) == 2)

    for label, override in (("no bundles", {"bundles": []}),
                            ("no finite", {"finite_licenses": []}),
                            ("neither", {"bundles": [], "finite_licenses": []})):
        d = dict(data, **override)
        b = build_slide(d)
        check(f"builds with {label}", len(b.getvalue()) > 10000)

    all_finite = [
        {"count": 142, "license_type": "Named User or Computer Based",
         "license_name": "DIAdem Professional with DAC"},
        {"count": 53, "license_type": "Named User or Computer Based",
         "license_name": "Circuit Design Suite"},
        {"count": 40, "license_type": "Named User or Computer Based",
         "license_name": "VeriStand Full"},
        {"count": 23, "license_type": "Concurrent",
         "license_name": "SystemLink Server Test Operations Suite Server"},
        {"count": 1, "license_type": "Concurrent",
         "license_name": "SystemLink Enterprise Test Operations Suite Server"},
        {"count": 165, "license_type": "Concurrent",
         "license_name": "SystemLink Test Operations Suite User"},
        {"count": 325, "license_type": "Concurrent",
         "license_name": "SystemLink Test Operations Suite Node"},
    ]
    prs_all = Presentation(build_slide(dict(data, finite_licenses=all_finite)))
    table_text = "\n".join(
        cell.text
        for slide in prs_all.slides
        for shape in slide.shapes
        if getattr(shape, "has_table", False)
        for row in shape.table.rows
        for cell in row.cells
    )
    check("finite table keeps all contract software",
          "SystemLink Test Operations Suite Node" in table_text)

    deck = build_deck([(data, ins), (dict(data, bundles=[]), None)])
    check("deck slide count", len(list(Presentation(deck).slides)) == 3)


def test_preview():
    from preview import generate_preview_html
    from slide_preview import preview_renderer_status

    html = generate_preview_html(_sample_data())
    for token in ("EA-15725", "Blue Origin", "<svg", "Peak machines",
                  "EA Platform Bundle", "TOTAL", "TOP VER."):
        check(f"preview contains {token!r}", token in html)
    status = preview_renderer_status()
    check("pptx preview status shape",
          {"platform", "quicklook", "libreoffice", "pdftoppm", "available"}
          <= set(status))


def test_profiles():
    import profiles as prof

    with tempfile.TemporaryDirectory() as tmp:
        orig = prof.PROFILES_DIR
        prof.PROFILES_DIR = Path(tmp)
        try:
            state = {"machine_text": "P\tN\tE\nQ1\t1\t2",
                     "locations_text": "", "versions_text": "",
                     "f_service_id": "EA-15725", "f_customer": "Blue Origin"}
            prof.save_profile("test", state,
                              [{"count": 5, "license_type": "t",
                                "license_name": "n"}],
                              ["B1"])
            check("profile listed", prof.list_profiles() == ["test"])
            p = prof.load_profile("test")
            check("profile texts round-trip",
                  p["texts"]["machine_text"] == state["machine_text"])
            check("profile fields round-trip",
                  p["fields"]["f_service_id"] == "EA-15725")
            check("profile finite round-trip",
                  p["finite_licenses"][0]["count"] == 5)
            check("profile bundles round-trip", p["bundles"] == ["B1"])
            check("profile name suggestion",
                  prof.suggest_name(state) == "EA-15725_Blue_Origin")
            check("profile delete", prof.delete_profile("test")
                  and prof.list_profiles() == [])
        finally:
            prof.PROFILES_DIR = orig


def test_windows_batch_files():
    root = Path(__file__).resolve().parents[1]
    start = (root / "Start EA Slide Builder.bat").read_text(encoding="utf-8")
    update = (root / "Update EA Slide Builder.bat").read_text(encoding="utf-8")
    build = (root / "packaging" / "build_exe.bat").read_text(encoding="utf-8")

    check("windows start uses setlocal", "setlocal" in start)
    check("windows start supports py launcher", "py -3" in start)
    check("windows start uses venv python",
          '.venv\\Scripts\\python.exe' in start)
    check("windows start is localhost-only",
          "--server.address=127.0.0.1" in start)
    check("windows update fast-forwards main",
          "git pull --ff-only origin main" in update)
    check("windows exe build uses venv python",
          '.venv\\Scripts\\python.exe' in build)
    check("windows exe build runs PyInstaller",
          "-m PyInstaller" in build)


def test_github_pages_lite_security():
    root = Path(__file__).resolve().parents[1]
    index = (root / "docs" / "index.html").read_text(encoding="utf-8")
    app_js = (root / "docs" / "assets" / "app.js").read_text(encoding="utf-8")
    styles = (root / "docs" / "assets" / "styles.css").read_text(encoding="utf-8")
    combined = "\n".join([index, app_js, styles])

    check("pages lite index exists", "<title>EA Slide Builder Lite</title>" in index)
    check("pages lite has CSP", "Content-Security-Policy" in index)
    check("pages lite limits network to self", "connect-src 'self'" in index)
    check("pages lite uses local OCR worker", "worker-src 'self'" in index)
    check("pages lite uses local script", 'src="assets/app.js"' in index)
    check("pages lite uses local stylesheet", 'href="assets/styles.css"' in index)
    check("pages lite has screenshot uploads", 'id="screenshotA" type="file"' in index)
    check("pages lite has PPTX generation", "downloadCurrentPptx" in app_js)
    check("pages lite has profile import/export", "currentProfilePayload" in app_js)
    check("pages lite has batch deck generation", "downloadBatchPptx" in app_js)
    check("pages lite has browser OCR", "Tesseract.createWorker" in app_js)
    check("pages lite disables OCR cache", "cacheMethod: 'none'" in app_js)
    check("pages lite has no external urls",
          "http://" not in combined and "https://" not in combined)
    check("pages lite has no browser storage",
          "localStorage" not in combined and "sessionStorage" not in combined)
    vendor = root / "docs" / "assets" / "vendor"
    check("pages lite vendors JSZip", (vendor / "jszip.min.js").is_file())
    check("pages lite vendors PptxGenJS", (vendor / "pptxgen.bundle.js").is_file())
    check("pages lite vendors Tesseract API",
          (vendor / "tesseract" / "tesseract.min.js").is_file())
    check("pages lite vendors Tesseract worker",
          (vendor / "tesseract" / "worker.min.js").is_file())
    check("pages lite vendors English OCR data",
          (vendor / "tesseract" / "lang" / "eng.traineddata.gz").is_file())
    core = vendor / "tesseract" / "core"
    for core_name in (
        "tesseract-core.wasm.js",
        "tesseract-core-lstm.wasm.js",
        "tesseract-core-simd.wasm.js",
        "tesseract-core-simd-lstm.wasm.js",
        "tesseract-core-relaxedsimd.wasm.js",
        "tesseract-core-relaxedsimd-lstm.wasm.js",
    ):
        check(f"pages lite vendors {core_name}", (core / core_name).is_file())


# --------------------------------------------------------------------------- #
# Full app via Streamlit AppTest
# --------------------------------------------------------------------------- #
def test_app():
    from streamlit.testing.v1 import AppTest

    app_path = str(Path(__file__).resolve().parents[1] / "app.py")
    at = AppTest.from_file(app_path, default_timeout=60).run()
    check("app runs with no input", not at.exception,
          str(at.exception[0].value) if at.exception else "")

    at.text_area(key="machine_text").set_value(
        "P\tNew\tExisting\nQ1\t10\t100\nQ2\t20\t120")
    at.text_input(key="f_flex_used").set_value("5000")
    at.text_input(key="f_flex_purchased").set_value("20000")
    at.run()
    check("app runs with inputs", not at.exception,
          str(at.exception[0].value) if at.exception else "")

    btns = [b for b in at.button if "Generate Slide" in b.label]
    check("generate button present", len(btns) == 1)
    if btns:
        btns[0].click().run()
        check("generate path clean", not at.exception,
              str(at.exception[0].value) if at.exception else "")
        check("success shown",
              any("Slide generated" in s.value for s in at.success))


def main() -> int:
    print("=" * 60)
    print("EA Slide Builder self-test")
    print("=" * 60)
    for fn in (test_parsers, test_screenshot_parsers, test_slide_builder,
               test_preview, test_profiles, test_windows_batch_files,
               test_github_pages_lite_security, test_app):
        print(f"\n--- {fn.__name__} ---")
        try:
            fn()
        except Exception as exc:  # a crashed section is a failure, not an abort
            check(f"{fn.__name__} crashed", False, repr(exc))
    print("\n" + "=" * 60)
    if FAILURES:
        print(f"RESULT: {len(FAILURES)} FAILURE(S): {', '.join(FAILURES)}")
        return 1
    print("RESULT: all checks passed")
    return 0


if __name__ == "__main__":
    sys.exit(main())
