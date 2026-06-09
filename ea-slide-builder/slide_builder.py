"""
Builds a single 16:9 .pptx slide matching the EA account summary layout.

Layout (matching reference design):
  Title (full width, top-left) + "Updated on" (top-right)
  3 columns below title:
    LEFT  : rounded info box (bullet points) + Bundle table + Finite Licenses table
    CENTER: Software Usage Trend chart + Software Usage Data table
    RIGHT : Top Site Locations + Version Usage + Training Credits + Technical Support
"""
from __future__ import annotations

from datetime import date
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

# ---------------------------------------------------------------------------
# Theme colours
# ---------------------------------------------------------------------------
DARK_GREEN   = RGBColor(0x01, 0x33, 0x24)
ACCENT_GREEN = RGBColor(0x18, 0xAF, 0x7C)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x00, 0x00, 0x00)
LIGHT_GREY   = RGBColor(0xF2, 0xF2, 0xF2)
MID_GREY     = RGBColor(0xD9, 0xD9, 0xD9)

# ---------------------------------------------------------------------------
# Slide & layout constants  (all Emu)
# ---------------------------------------------------------------------------
SLIDE_W = Emu(9144000)   # 10 in
SLIDE_H = Emu(5143500)   # 5.625 in  (16:9)

MARGIN  = Emu(200000)
GAP     = Emu(100000)

TITLE_H  = Emu(480000)
COL_TOP  = Emu(int(MARGIN) + int(TITLE_H) + int(GAP))
COL_H    = Emu(int(SLIDE_H) - int(COL_TOP) - int(MARGIN))
COL_W    = Emu((int(SLIDE_W) - int(MARGIN)*2 - int(GAP)*2) // 3)

LEFT_X   = MARGIN
CENTER_X = Emu(int(LEFT_X)   + int(COL_W) + int(GAP))
RIGHT_X  = Emu(int(CENTER_X) + int(COL_W) + int(GAP))


# ---------------------------------------------------------------------------
# Row height helpers
# ---------------------------------------------------------------------------
def _rh(fs: float) -> int:
    return int(Pt(fs + 3).emu) + 45000


def _auto_fs(n_data: int, avail_h: int, max_fs: float = 8.0, min_fs: float = 5.0) -> float:
    fs = max_fs
    while fs > min_fs:
        if _rh(fs) * (n_data + 1) <= avail_h:   # +1 for header
            return fs
        fs -= 0.5
    return min_fs


# ---------------------------------------------------------------------------
# Drawing primitives
# ---------------------------------------------------------------------------

def _txb(slide, left, top, width, height, text,
         bold=False, size=10.0, color=BLACK, align=PP_ALIGN.LEFT, bg=None):
    s = slide.shapes.add_textbox(left, top, width, height)
    if bg:
        s.fill.solid()
        s.fill.fore_color.rgb = bg
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.bold  = bold
    r.font.size  = Pt(size)
    r.font.color.rgb = color
    return s


def _hdr(slide, left, top, width, title: str, fs: float = 8.0) -> Emu:
    h = Emu(_rh(fs) + 20000)
    _txb(slide, left, top, width, h, title,
         bold=True, size=fs, color=WHITE, bg=DARK_GREEN, align=PP_ALIGN.CENTER)
    return Emu(int(top) + int(h))


def _table(slide, left, top, width, rows, col_widths, fs: float = 7.5):
    """Draw a table; rows[0] is the header. Returns bottom Y."""
    nr = len(rows)
    nc = len(rows[0])
    rh = _rh(fs)
    tbl = slide.shapes.add_table(nr, nc, left, top, width, rh * nr).table

    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ri, row in enumerate(rows):
        is_hdr = ri == 0
        bg_clr = DARK_GREEN if is_hdr else (LIGHT_GREY if ri % 2 == 0 else WHITE)
        fg_clr = WHITE if is_hdr else BLACK
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.size  = Pt(fs)
            p.font.bold  = is_hdr
            p.font.color.rgb = fg_clr
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_clr

    return Emu(int(top) + rh * nr)


def _rounded_box_bullets(slide, left, top, width, height, bullets: list[str]):
    """Rounded box with ACCENT_GREEN border and a bullet list inside."""
    shape = slide.shapes.add_shape(5, left, top, width, height)
    try:
        if len(shape.adjustments) > 0:
            shape.adjustments[0] = 0.05
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = ACCENT_GREEN
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True

    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        r = p.add_run()
        r.text = u"•  " + text
        r.font.size  = Pt(8.5)
        r.font.color.rgb = DARK_GREEN


def _line_chart(slide, left, top, width, height, periods, totals):
    cd = ChartData()
    cd.categories = periods
    cd.add_series("Sessions", [float(t) for t in totals])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, left, top, width, height, cd
    ).chart
    chart.series[0].format.line.color.rgb = ACCENT_GREEN
    chart.series[0].format.line.width = Pt(2)
    chart.has_legend = False
    try:
        chart.font.size = Pt(7)
    except Exception:
        pass
    for attr in ("category_axis", "value_axis"):
        try:
            getattr(chart, attr).tick_labels.font.size = Pt(6)
        except Exception:
            pass


# ---------------------------------------------------------------------------
# Main builder
# ---------------------------------------------------------------------------

def build_slide(data: dict[str, Any], contract: dict[str, Any], out_path: str | Path) -> Path:
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # ------------------------------------------------------------------
    # Title  (left)  +  "Updated on" (right)
    # ------------------------------------------------------------------
    ea_num   = contract.get("ea_number", "EA-XXXXXX")
    customer = contract.get("customer_name", "Customer")
    _txb(slide, MARGIN, MARGIN,
         Emu(int(SLIDE_W) - int(MARGIN)*2 - 600000), TITLE_H,
         f"{ea_num} - {customer}",
         bold=True, size=16, color=DARK_GREEN)

    today_str = f"Updated on: {date.today().strftime('%-m/%-d/%Y')}"
    _txb(slide, Emu(int(SLIDE_W) - 700000 - int(MARGIN)), MARGIN,
         Emu(700000), TITLE_H,
         today_str, size=7, color=RGBColor(0x80, 0x80, 0x80), align=PP_ALIGN.RIGHT)

    # ==================================================================
    # LEFT COLUMN
    # ==================================================================
    y = COL_TOP

    # --- Rounded info box (bullet points) ---
    end_date = contract.get("ea_end_date", "N/A")
    term     = contract.get("term_duration", "N/A")
    scope    = contract.get("contract_scope", "N/A")
    phase    = contract.get("phase", "N/A")

    bullets = [
        f"EA End Date: {end_date}",
        f"Term Duration: {term}",
        f"Contract Scope: {scope}",
        f"Phase: {phase}",
    ]
    box_h = Emu(_rh(8.5) * len(bullets) + 80000)
    _rounded_box_bullets(slide, LEFT_X, y, COL_W, box_h, bullets)
    y = Emu(int(y) + int(box_h) + int(GAP))

    # --- Bundle Information ---
    bundles = contract.get("bundles", [])
    y = _hdr(slide, LEFT_X, y, COL_W, "Bundle Information")
    b_rows = [["Bundle Name", "Type"]]
    for b in bundles:
        b_rows.append([b.get("name", ""), b.get("type", "")])
    if not bundles:
        b_rows.append(["No bundles defined", ""])
    cw_b = [int(COL_W * 0.65), int(COL_W * 0.35)]
    y = _table(slide, LEFT_X, y, COL_W, b_rows, cw_b, fs=7.5)
    y = Emu(int(y) + int(GAP))

    # --- Finite Quantity Licenses (dynamic) ---
    fql = contract.get("finite_quantity_licenses", [])
    y = _hdr(slide, LEFT_X, y, COL_W, "Finite Quantity NI SW Licenses")

    avail_lic = int(COL_H) - (int(y) - int(COL_TOP)) - int(MARGIN)
    n_lic     = max(1, len(fql))
    lic_fs    = _auto_fs(n_lic, avail_lic, max_fs=7.5, min_fs=5.0)

    l_rows = [["Qty", "Product", "Type"]]
    for lic in fql:
        l_rows.append([
            str(lic.get("count", "")),
            lic.get("name", ""),
            lic.get("type", ""),
        ])
    if not fql:
        l_rows.append(["--", "No licenses defined", ""])
    cw_l = [int(COL_W * 0.13), int(COL_W * 0.52), int(COL_W * 0.35)]
    _table(slide, LEFT_X, y, COL_W, l_rows, cw_l, fs=lic_fs)

    # ==================================================================
    # CENTER COLUMN
    # ==================================================================
    yc = COL_TOP
    machines = data.get("machines")

    yc = _hdr(slide, CENTER_X, yc, COL_W, "Software Usage Trend")
    chart_h = Emu(int(COL_H) * 55 // 100)

    if machines and len(machines["periods"]) >= 2:
        _line_chart(slide, CENTER_X, yc, COL_W, chart_h,
                    machines["periods"], machines["totals"])
    else:
        _txb(slide, CENTER_X, yc, COL_W, chart_h,
             "Insufficient period data for trend chart (need 2+ quarters/months)",
             size=7, color=BLACK)
    yc = Emu(int(yc) + int(chart_h) + int(GAP))

    yc = _hdr(slide, CENTER_X, yc, COL_W, "Software Usage Data")
    if machines:
        max_s = f"{machines['max_sessions']:,} ({machines['max_period']})"
        min_s = f"{machines['min_sessions']:,} ({machines['min_period']})"
        avg_c = f"{machines['avg_pct_change']:+.0f}%"
    else:
        max_s = min_s = avg_c = "N/A"

    usage_rows = [
        ["Metric", "Value"],
        ["Overall Software Usage Maximum Amount (Machines)", max_s],
        ["Overall Software Usage Minimum Amount (Machines)", min_s],
        ["Average Percent of Increase per Quarter %", avg_c],
    ]
    cw_u = [int(COL_W * 0.62), int(COL_W * 0.38)]
    _table(slide, CENTER_X, yc, COL_W, usage_rows, cw_u, fs=7.0)

    # ==================================================================
    # RIGHT COLUMN
    # ==================================================================
    yr = COL_TOP

    # --- Top Site Locations (Country | Region | City | Machines) ---
    yr = _hdr(slide, RIGHT_X, yr, COL_W, "Top Site Locations")
    cities = data.get("cities")
    if cities and cities["top_locations"]:
        loc_rows = [["Country", "State/Region", "City", "No. of Machines"]]
        for loc in cities["top_locations"]:
            loc_rows.append([
                loc.get("country", ""),
                loc.get("region", ""),
                loc.get("city", ""),
                f"{loc['usage']:,}",
            ])
    else:
        loc_rows = [["Country", "State/Region", "City", "No. of Machines"],
                    ["No data", "", "", "--"]]
    cw_s = [int(COL_W*0.28), int(COL_W*0.24), int(COL_W*0.26), int(COL_W*0.22)]
    yr = _table(slide, RIGHT_X, yr, COL_W, loc_rows, cw_s, fs=7.0)
    yr = Emu(int(yr) + int(GAP))

    # --- Version Usage (top N products) ---
    yr = _hdr(slide, RIGHT_X, yr, COL_W, "Version Usage")
    software = data.get("software")
    if software and software.get("top_version_rows"):
        ver_rows = [["Product", "Version", "Version Usage", "% Usage"]]
        for row in software["top_version_rows"]:
            ver_rows.append([
                row["product"][:20],
                row["top_version"][:10],
                f"{row['top_version_usage']:,}",
                f"{row['top_version_pct']}%",
            ])
    elif software and software["products"]:
        ver_rows = [["Product", "Version", "Version Usage", "% Usage"]]
        for p in software["products"]:
            ver_rows.append([p["product"][:20], p["top_version"][:10],
                             f"{p['usage']:,}", ""])
    else:
        ver_rows = [["Product", "Version", "Version Usage", "% Usage"],
                    ["No data", "--", "--", "--"]]

    # auto-size to fit remaining right-column space
    remaining_r = int(COL_H) - (int(yr) - int(COL_TOP)) - int(GAP)*3 - _rh(7.5)*5
    ver_fs = _auto_fs(len(ver_rows)-1, max(remaining_r, _rh(7.5)*3), max_fs=7.5, min_fs=5.5)
    cw_v = [int(COL_W*0.35), int(COL_W*0.20), int(COL_W*0.24), int(COL_W*0.21)]
    yr = _table(slide, RIGHT_X, yr, COL_W, ver_rows, cw_v, fs=ver_fs)
    yr = Emu(int(yr) + int(GAP))

    # --- Training Credit Usage ---
    yr = _hdr(slide, RIGHT_X, yr, COL_W, "Training Credit Usage")
    tc_total = int(contract.get("training_credits_total") or 0)
    tc_used  = int(contract.get("training_credits_used")  or 0)
    tc_pct   = f"{round(tc_used / tc_total * 100, 0):.0f}%" if tc_total else "N/A"
    tc_rows = [
        ["Purchased", "Used", "% Used"],
        [str(tc_total) if tc_total else "N/A",
         str(tc_used)  if tc_total else "N/A",
         tc_pct],
    ]
    cw_t = [int(COL_W*0.34), int(COL_W*0.33), int(COL_W*0.33)]
    yr = _table(slide, RIGHT_X, yr, COL_W, tc_rows, cw_t, fs=7.5)
    yr = Emu(int(yr) + int(GAP))

    # --- Technical Support ---
    yr = _hdr(slide, RIGHT_X, yr, COL_W, "Technical Support")
    support = contract.get("technical_support_level", "Standard")
    ts_rows = [
        ["Support Level", "Users"],
        [support, "All Users"],
    ]
    cw_ts = [int(COL_W*0.55), int(COL_W*0.45)]
    _table(slide, RIGHT_X, yr, COL_W, ts_rows, cw_ts, fs=7.5)

    prs.save(out_path)
    return out_path
