"""slide_builder.py

Builds the single 16:9 EA one-slider with python-pptx, matching the reference
visual design:

  - full-width dark-green header band (EA number + customer, "Updated <date>")
  - three columns of white rounded "cards" with thin gray borders
  - LEFT:   Contract Details, Bundle Information, NI SW Licenses (finite qty)
  - CENTER: Software Usage Trend (native line chart), Software Usage Data
            (peak/min stat callouts + avg-increase strip)
  - RIGHT:  Top Site Locations, Version Usage, Training Credit Usage,
            Technical Support (solid dark-green card)

Colors: dark green #013324, accent green #18af7c, light-green tint for
alternating rows / the avg-increase strip. Serif headline font, sans body.

All rendering is local (python-pptx writes the .pptx in memory). No network.
"""

from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------- #
# Theme
# --------------------------------------------------------------------------- #
DARK_GREEN = RGBColor(0x01, 0x33, 0x24)
ACCENT_GREEN = RGBColor(0x18, 0xAF, 0x7C)
LIGHT_TINT = RGBColor(0xEC, 0xF7, 0xF2)   # very light green, alternating rows
CARD_BORDER = RGBColor(0xDD, 0xDD, 0xDD)  # thin light-gray card border
GRAY_TEXT = RGBColor(0x6E, 0x6E, 0x6E)
MUTED_GREEN = RGBColor(0x8F, 0xC9, 0xB6)  # muted green for header label
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SERIF = "Georgia"        # headline numbers + title
SANS = "Calibri"         # body + table text

# Slide geometry (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --------------------------------------------------------------------------- #
# Low-level helpers
# --------------------------------------------------------------------------- #
def _set_letter_spacing(run, centipoints: int = 150) -> None:
    """Apply letter spacing to a run (spc is in 1/100 pt)."""
    run.font._rPr.set("spc", str(centipoints))


def _no_autofit(tf) -> None:
    tf.word_wrap = True


def _set_text(cell_or_tf, text, *, size=10, bold=False, color=DARK_GREEN,
              font=SANS, align=PP_ALIGN.LEFT, anchor=None, spacing=None):
    """Write a single styled paragraph into a text frame (or table cell)."""
    tf = cell_or_tf.text_frame if hasattr(cell_or_tf, "text_frame") else cell_or_tf
    tf.clear()
    _no_autofit(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    if spacing:
        _set_letter_spacing(run, spacing)
    if anchor is not None:
        tf.vertical_anchor = anchor
    return tf


def _rect(slide, x, y, w, h, *, fill=None, line=None, line_w=Pt(0.75),
          rounded=True):
    """Add a (rounded) rectangle and return the shape."""
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shp, x, y, w, h)
    shape.shadow.inherit = False
    if rounded:
        # Gentle corner radius.
        try:
            shape.adjustments[0] = 0.06
        except (IndexError, KeyError):
            pass
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = line_w
    return shape


def _card(slide, x, y, w, h, title, *, title_color=ACCENT_GREEN,
          body_color=DARK_GREEN):
    """Draw a white rounded card with a thin gray border and a small uppercase
    letter-spaced section title. Returns the inner content rect (x, y, w, h)."""
    fill = body_color if title_color == WHITE else WHITE
    _rect(slide, x, y, w, h, fill=fill, line=CARD_BORDER, line_w=Pt(0.75))

    pad = Inches(0.14)
    title_h = Inches(0.26)
    tb = slide.shapes.add_textbox(x + pad, y + Inches(0.08), w - 2 * pad, title_h)
    _set_text(tb, title.upper(), size=8.5, bold=True, color=title_color,
              font=SANS, spacing=180)

    inner_x = x + pad
    inner_y = y + Inches(0.08) + title_h
    inner_w = w - 2 * pad
    inner_h = h - (Inches(0.08) + title_h) - pad
    return inner_x, inner_y, inner_w, inner_h


def _empty_note(slide, x, y, w, h, text):
    """A tidy, centered muted placeholder for an optional card with no data."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    _set_text(tb, text, size=9, color=GRAY_TEXT, font=SANS,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _style_table(table):
    """Strip python-pptx's default banded style so our explicit fills win."""
    table.first_row = False
    table.horz_banding = False
    # Set a plain, light table style if available; otherwise leave fills to us.


def _cell(cell, text, *, size=8.5, bold=False, color=DARK_GREEN, font=SANS,
          align=PP_ALIGN.LEFT, fill=None):
    """Style a single table cell."""
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.01)
    cell.margin_bottom = Inches(0.01)
    _set_text(cell, text, size=size, bold=bold, color=color, font=font,
              align=align)


def _set_table_row_heights(table, total_h, n_rows: int) -> None:
    """Force rows to share the available table height."""
    if n_rows <= 0:
        return
    row_h = Emu(max(1, int(total_h) // n_rows))
    for row in table.rows:
        row.height = row_h


def _adaptive_body_size(total_h, n_rows: int, *, base=8.0, minimum=5.4) -> float:
    """Choose a table font size from the available row height."""
    if n_rows <= 0:
        return base
    row_h_inches = int(total_h) / 914400 / n_rows
    return max(minimum, min(base, row_h_inches * 25))


# --------------------------------------------------------------------------- #
# Header band
# --------------------------------------------------------------------------- #
def _header(slide, data):
    band_h = Inches(0.92)
    _rect(slide, 0, 0, SLIDE_W, band_h, fill=DARK_GREEN, rounded=False)

    pad = Inches(0.35)
    # Small uppercase muted-green label.
    tb = slide.shapes.add_textbox(pad, Inches(0.12), Inches(8), Inches(0.24))
    _set_text(tb, "ENTERPRISE AGREEMENT", size=8.5, bold=True,
              color=MUTED_GREEN, font=SANS, spacing=200)

    # EA number · customer, bold white serif.
    sid = data.get("service_id") or "EA-—"
    cust = data.get("customer") or "Customer"
    tb2 = slide.shapes.add_textbox(pad, Inches(0.36), Inches(10), Inches(0.5))
    tf = tb2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for txt in (sid, "  ·  ", cust):
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.name = SERIF
        r.font.color.rgb = WHITE

    # Right-aligned "Updated <date>".
    updated = data.get("updated_date", "")
    tb3 = slide.shapes.add_textbox(SLIDE_W - Inches(3.5), Inches(0.30),
                                   Inches(3.15), Inches(0.3))
    _set_text(tb3, f"Updated {updated}" if updated else "", size=9,
              color=MUTED_GREEN, font=SANS, align=PP_ALIGN.RIGHT)


# --------------------------------------------------------------------------- #
# LEFT column cards
# --------------------------------------------------------------------------- #
def _contract_details_card(slide, x, y, w, h, data):
    ix, iy, iw, ih = _card(slide, x, y, w, h, "Contract Details")
    rows = [
        ("EA End Date", data.get("ea_end_date", "")),
        ("Term Duration", data.get("ep_term", "")),
        ("Contract Scope", data.get("contract_scope", "")),
        ("Phase", data.get("phase", "")),
    ]
    row_h = ih / len(rows)
    for i, (label, value) in enumerate(rows):
        ry = iy + Emu(int(row_h) * i)
        lb = slide.shapes.add_textbox(ix, ry, iw * 0.5, Emu(int(row_h)))
        _set_text(lb, label, size=9, color=GRAY_TEXT, font=SANS,
                  anchor=MSO_ANCHOR.MIDDLE)
        vb = slide.shapes.add_textbox(ix + iw * 0.42, ry, iw * 0.58,
                                      Emu(int(row_h)))
        _set_text(vb, value or "—", size=9.5, bold=True, color=DARK_GREEN,
                  font=SANS, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _bundle_card(slide, x, y, w, h, data):
    ix, iy, iw, ih = _card(slide, x, y, w, h, "Bundle Information")
    bundles = data.get("bundles", []) or []
    if not bundles:
        _empty_note(slide, ix, iy, iw, ih, "No bundles provided")
        return
    pill_h = Inches(0.34)
    gap = Inches(0.08)
    max_pills = max(1, int(ih / (pill_h + gap)))
    for i, name in enumerate(bundles[:max_pills]):
        py = iy + Emu(int(pill_h + gap) * i)
        pill = _rect(slide, ix, py, iw, pill_h, fill=WHITE, line=ACCENT_GREEN,
                     line_w=Pt(1.0))
        _set_text(pill, name, size=9, bold=True, color=DARK_GREEN, font=SANS,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _finite_card(slide, x, y, w, h, data):
    ix, iy, iw, ih = _card(slide, x, y, w, h, "NI SW Licenses (Finite Qty)")
    all_rows = data.get("finite_licenses", []) or []
    if not all_rows:
        _empty_note(slide, ix, iy, iw, ih, "No finite licenses provided")
        return
    rows = all_rows
    headers = ["QTY", "LICENSE", "TYPE"]
    n_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n_rows, 3, ix, iy, iw, ih)
    table = tbl_shape.table
    _style_table(table)
    _set_table_row_heights(table, ih, n_rows)
    table.columns[0].width = Emu(int(iw * 0.14))
    table.columns[1].width = Emu(int(iw * 0.56))
    table.columns[2].width = Emu(int(iw * 0.30))

    body_size = _adaptive_body_size(ih, n_rows, base=8.0, minimum=5.2)
    header_size = max(5.3, min(7.4, body_size - 0.2))

    for c, htext in enumerate(headers):
        _cell(table.cell(0, c), htext, size=header_size, bold=True, color=WHITE,
              font=SANS, fill=DARK_GREEN,
              align=PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT)

    for r, row in enumerate(rows, start=1):
        shade = LIGHT_TINT if r % 2 == 0 else WHITE
        _cell(table.cell(r, 0), row.get("count", ""), size=body_size + 0.3, bold=True,
              color=ACCENT_GREEN, font=SANS, align=PP_ALIGN.CENTER, fill=shade)
        _cell(table.cell(r, 1), row.get("license_name", ""), size=body_size,
              color=DARK_GREEN, font=SANS, fill=shade)
        _cell(table.cell(r, 2), row.get("license_type", ""),
              size=max(5.0, body_size - 0.8),
              color=GRAY_TEXT, font=SANS, fill=shade)


# --------------------------------------------------------------------------- #
# CENTER column cards
# --------------------------------------------------------------------------- #
def _trend_card(slide, x, y, w, h, data):
    ix, iy, iw, ih = _card(slide, x, y, w, h, "Software Usage Trend")
    machine = data.get("machine", {})
    df = machine.get("df")
    if df is None or getattr(df, "empty", True):
        nb = slide.shapes.add_textbox(ix, iy, iw, Inches(0.3))
        _set_text(nb, "No machine-count data", size=9, color=GRAY_TEXT)
        return

    periods = [str(p) for p in df["period"].tolist()]
    totals = [int(t) for t in df["total"].tolist()]

    chart_data = CategoryChartData()
    chart_data.categories = periods
    chart_data.add_series("Total Machines", totals)

    # Least-squares linear trend, drawn as a subtle dashed guide line.
    n = len(totals)
    trend = None
    if n >= 3:
        xs = list(range(n))
        mean_x = sum(xs) / n
        mean_y = sum(totals) / n
        denom = sum((x - mean_x) ** 2 for x in xs) or 1
        slope = sum((x - mean_x) * (y - mean_y)
                    for x, y in zip(xs, totals)) / denom
        trend = [round(mean_y + slope * (x - mean_x), 1) for x in xs]
        chart_data.add_series("Trend", trend)

    # Peak marker: a one-point series so the max stands out on the line.
    peak_idx = totals.index(max(totals))
    chart_data.add_series(
        "Peak", [t if i == peak_idx else None for i, t in enumerate(totals)])

    gframe = slide.shapes.add_chart(XL_CHART_TYPE.LINE, ix, iy, iw, ih,
                                    chart_data)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = False

    series = chart.plots[0].series[0]
    series.smooth = False
    series.format.line.color.rgb = ACCENT_GREEN
    series.format.line.width = Pt(2.5)

    all_series = list(chart.plots[0].series)
    if trend is not None:
        tser = all_series[1]
        tser.smooth = False
        tser.format.line.color.rgb = RGBColor(0xB5, 0xD6, 0xC9)
        tser.format.line.width = Pt(1.25)
        tser.format.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    pser = all_series[-1]
    pser.smooth = False
    pser.format.line.fill.background()  # no connecting line, marker only
    pser.marker.style = XL_MARKER_STYLE.CIRCLE
    pser.marker.size = 8
    pser.marker.format.fill.solid()
    pser.marker.format.fill.fore_color.rgb = ACCENT_GREEN
    pser.marker.format.line.color.rgb = WHITE

    val_axis = chart.value_axis
    val_axis.has_major_gridlines = False
    val_axis.has_minor_gridlines = False
    val_axis.tick_labels.font.size = Pt(7)
    val_axis.tick_labels.font.name = SANS

    cat_axis = chart.category_axis
    cat_axis.has_major_gridlines = False
    cat_axis.tick_labels.font.size = Pt(7)
    cat_axis.tick_labels.font.name = SANS


def _usage_data_card(slide, x, y, w, h, data):
    ix, iy, iw, ih = _card(slide, x, y, w, h, "Software Usage Data")
    stats = data.get("machine", {}).get("stats", {})
    peak = stats.get("max_total", 0)
    peak_period = stats.get("max_period", "—")
    low = stats.get("min_total", 0)
    low_period = stats.get("min_period", "—")
    avg = stats.get("avg_pct_change", 0.0)

    gap = Inches(0.12)
    callout_w = (iw - gap) / 2
    callout_h = ih * 0.62

    # Left callout: peak (accent border, accent serif number).
    _stat_callout(slide, ix, iy, Emu(int(callout_w)), Emu(int(callout_h)),
                  peak, "Peak machines", peak_period,
                  number_color=ACCENT_GREEN, border=ACCENT_GREEN)
    # Right callout: min (no border emphasis, dark serif number).
    _stat_callout(slide, ix + Emu(int(callout_w + gap)), iy,
                  Emu(int(callout_w)), Emu(int(callout_h)),
                  low, "Min machines", low_period,
                  number_color=DARK_GREEN, border=CARD_BORDER)

    # Avg-increase strip (full width, light-green tint).
    strip_y = iy + Emu(int(callout_h)) + Inches(0.1)
    strip_h = ih - Emu(int(callout_h)) - Inches(0.1)
    _rect(slide, ix, strip_y, iw, strip_h, fill=LIGHT_TINT, line=None)
    lbl = slide.shapes.add_textbox(ix + Inches(0.12), strip_y, iw * 0.6, strip_h)
    _set_text(lbl, "Avg quarterly increase", size=9.5, color=DARK_GREEN,
              font=SANS, anchor=MSO_ANCHOR.MIDDLE)
    val = slide.shapes.add_textbox(ix + iw * 0.55, strip_y, iw * 0.45 - Inches(0.12),
                                   strip_h)
    _set_text(val, f"{avg:+.1f}%", size=20, bold=True, color=ACCENT_GREEN,
              font=SERIF, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _stat_callout(slide, x, y, w, h, number, label, period, *, number_color,
                  border):
    _rect(slide, x, y, w, h, fill=WHITE, line=border, line_w=Pt(1.25))
    tf = slide.shapes.add_textbox(x, y + Inches(0.08), w, h - Inches(0.12)).text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Big number
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"{int(number):,}"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.name = SERIF
    r.font.color.rgb = number_color
    # Label
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(9)
    r2.font.name = SANS
    r2.font.color.rgb = GRAY_TEXT
    # Period
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = str(period)
    r3.font.size = Pt(8)
    r3.font.name = SANS
    r3.font.color.rgb = DARK_GREEN


# --------------------------------------------------------------------------- #
# RIGHT column cards
# --------------------------------------------------------------------------- #
def _locations_card(slide, x, y, w, h, data):
    ix, iy, iw, ih = _card(slide, x, y, w, h, "Top Site Locations")
    rows = (data.get("locations_top5", []) or [])[:5]
    if not rows:
        _empty_note(slide, ix, iy, iw, ih, "No location data")
        return
    headers = ["STATE", "CITY", "MACHINES"]
    table = slide.shapes.add_table(len(rows) + 1, 3, ix, iy, iw, ih).table
    _style_table(table)
    _set_table_row_heights(table, ih, len(rows) + 1)
    table.columns[0].width = Emu(int(iw * 0.22))
    table.columns[1].width = Emu(int(iw * 0.48))
    table.columns[2].width = Emu(int(iw * 0.30))
    body_size = _adaptive_body_size(ih, len(rows) + 1, base=8.4, minimum=6.4)

    for c, htext in enumerate(headers):
        _cell(table.cell(0, c), htext, size=max(6.0, body_size - 0.8),
              bold=True, color=WHITE,
              font=SANS, fill=DARK_GREEN,
              align=PP_ALIGN.RIGHT if c == 2 else PP_ALIGN.LEFT)
    for r, row in enumerate(rows, start=1):
        shade = LIGHT_TINT if r % 2 == 0 else WHITE
        _cell(table.cell(r, 0), row.get("state", ""), size=body_size,
              color=DARK_GREEN, fill=shade)
        _cell(table.cell(r, 1), row.get("city") or row.get("location", ""),
              size=body_size, color=DARK_GREEN, fill=shade)
        _cell(table.cell(r, 2), f"{int(row.get('count', 0)):,}", size=body_size + 0.2,
              bold=True, color=ACCENT_GREEN, align=PP_ALIGN.RIGHT, fill=shade)


def _version_card(slide, x, y, w, h, data):
    ix, iy, iw, ih = _card(slide, x, y, w, h, "Version Usage")
    rows = (data.get("versions_top5", []) or [])[:5]
    if not rows:
        _empty_note(slide, ix, iy, iw, ih, "No version data")
        return
    headers = ["PRODUCT", "TOTAL", "TOP VER.", "%"]
    table = slide.shapes.add_table(len(rows) + 1, 4, ix, iy, iw, ih).table
    _style_table(table)
    _set_table_row_heights(table, ih, len(rows) + 1)
    table.columns[0].width = Emu(int(iw * 0.38))
    table.columns[1].width = Emu(int(iw * 0.19))
    table.columns[2].width = Emu(int(iw * 0.27))
    table.columns[3].width = Emu(int(iw * 0.16))
    body_size = _adaptive_body_size(ih, len(rows) + 1, base=7.6, minimum=5.9)

    for c, htext in enumerate(headers):
        _cell(table.cell(0, c), htext, size=max(5.6, body_size - 0.7),
              bold=True, color=WHITE,
              font=SANS, fill=DARK_GREEN,
              align=PP_ALIGN.RIGHT if c in (1, 3) else PP_ALIGN.LEFT)
    for r, row in enumerate(rows, start=1):
        shade = LIGHT_TINT if r % 2 == 0 else WHITE
        _cell(table.cell(r, 0), row.get("product", ""), size=body_size,
              color=DARK_GREEN, fill=shade)
        _cell(table.cell(r, 1), f"{int(row.get('product_total', row.get('users', 0))):,}",
              size=body_size, bold=True, color=DARK_GREEN,
              align=PP_ALIGN.RIGHT, fill=shade)
        _cell(table.cell(r, 2), row.get("version", ""), size=body_size,
              color=DARK_GREEN, align=PP_ALIGN.RIGHT, fill=shade)
        _cell(table.cell(r, 3), f"{int(row.get('pct', 0))}%", size=body_size + 0.2,
              bold=True, color=ACCENT_GREEN, align=PP_ALIGN.RIGHT, fill=shade)


def _credits_card(slide, x, y, w, h, data):
    ix, iy, iw, ih = _card(slide, x, y, w, h, "Training Credit Usage")
    credits = data.get("credits", {}) or {}
    purchased = credits.get("purchased", "—")
    used = credits.get("used", "—")
    pct = credits.get("pct_used", "—")
    cols = [
        ("Purchased", purchased, DARK_GREEN),
        ("Used", used, DARK_GREEN),
        ("Utilized", f"{pct}%" if pct not in ("", "—", None) else "—",
         ACCENT_GREEN),
    ]
    col_w = iw / 3
    for i, (label, value, color) in enumerate(cols):
        cx = ix + Emu(int(col_w) * i)
        lb = slide.shapes.add_textbox(cx, iy, Emu(int(col_w)), Inches(0.24))
        _set_text(lb, label, size=8, color=GRAY_TEXT, font=SANS,
                  align=PP_ALIGN.CENTER)
        vb = slide.shapes.add_textbox(cx, iy + Inches(0.26), Emu(int(col_w)),
                                      ih - Inches(0.26))
        _set_text(vb, _fmt(value), size=19, bold=True, color=color, font=SERIF,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


def _support_card(slide, x, y, w, h, data):
    # Solid dark-green card.
    _rect(slide, x, y, w, h, fill=DARK_GREEN, line=DARK_GREEN)
    pad = Inches(0.14)
    title_tb = slide.shapes.add_textbox(x + pad, y + Inches(0.08), w - 2 * pad,
                                        Inches(0.24))
    _set_text(title_tb, "TECHNICAL SUPPORT", size=8.5, bold=True,
              color=MUTED_GREEN, font=SANS, spacing=180)

    support = data.get("support", {}) or {}
    tier = support.get("tier", "—")
    scope = support.get("scope", "")
    body = slide.shapes.add_textbox(x + pad, y + Inches(0.34), w - 2 * pad,
                                    h - Inches(0.42))
    tf = body.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = tier
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.name = SANS
    r.font.color.rgb = WHITE
    if scope:
        r2 = p.add_run()
        r2.text = f"   {scope}"
        r2.font.size = Pt(10)
        r2.font.name = SANS
        r2.font.color.rgb = MUTED_GREEN


# --------------------------------------------------------------------------- #
# Misc
# --------------------------------------------------------------------------- #
def _fmt(value):
    try:
        return f"{int(str(value).replace(',', '')):,}"
    except (ValueError, TypeError):
        return str(value)


# --------------------------------------------------------------------------- #
# Insights slide
# --------------------------------------------------------------------------- #
_PRIORITY_LABELS = {1: "HIGH", 2: "MEDIUM", 3: "GOOD", 4: "CONTEXT"}
_PRIORITY_COLORS = {
    1: RGBColor(0xC0, 0x3A, 0x2B),   # muted red
    2: RGBColor(0xC8, 0x8A, 0x1E),   # amber
    3: ACCENT_GREEN,
    4: DARK_GREEN,
}


def _add_insights_slide(prs, data: dict, insights: list[dict]) -> None:
    """A second slide listing the prioritized CSM insights."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE, rounded=False)
    bg.line.fill.background()

    # Header band, matching the main slide.
    band_h = Inches(0.92)
    _rect(slide, 0, 0, SLIDE_W, band_h, fill=DARK_GREEN, rounded=False)
    pad = Inches(0.35)
    tb = slide.shapes.add_textbox(pad, Inches(0.12), Inches(8), Inches(0.24))
    _set_text(tb, "CSM INSIGHTS", size=8.5, bold=True, color=MUTED_GREEN,
              font=SANS, spacing=200)
    sid = data.get("service_id") or "EA-—"
    cust = data.get("customer") or "Customer"
    tb2 = slide.shapes.add_textbox(pad, Inches(0.36), Inches(10), Inches(0.5))
    tf = tb2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    for txt in (sid, "  ·  ", cust):
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.name = SERIF
        r.font.color.rgb = WHITE

    # Insight rows: priority chip + category + text, alternating tint.
    x = Inches(0.35)
    w = SLIDE_W - Inches(0.7)
    y = Inches(1.15)
    row_h = Inches(0.92)
    gap = Inches(0.1)
    for i, ins in enumerate(insights[:6]):
        fill = LIGHT_TINT if i % 2 == 0 else WHITE
        _rect(slide, x, y, w, row_h, fill=fill, line=CARD_BORDER,
              line_w=Pt(0.75))
        prio = int(ins.get("priority", 4))
        chip_w = Inches(1.0)
        chip = _rect(slide, x + Inches(0.12), y + Inches(0.3), chip_w,
                     Inches(0.32), fill=_PRIORITY_COLORS.get(prio, DARK_GREEN),
                     line=None)
        _set_text(chip, _PRIORITY_LABELS.get(prio, ""), size=8.5, bold=True,
                  color=WHITE, font=SANS, align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE, spacing=120)
        body = slide.shapes.add_textbox(x + Inches(1.3), y + Inches(0.06),
                                        w - Inches(1.45), row_h - Inches(0.12))
        btf = body.text_frame
        btf.word_wrap = True
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = btf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = str(ins.get("category", "")).upper()
        r1.font.size = Pt(8.5)
        r1.font.bold = True
        r1.font.name = SANS
        r1.font.color.rgb = ACCENT_GREEN
        p2 = btf.add_paragraph()
        r2 = p2.add_run()
        r2.text = str(ins.get("text", ""))
        r2.font.size = Pt(9.5)
        r2.font.name = SANS
        r2.font.color.rgb = DARK_GREEN
        y = Emu(int(y) + int(row_h) + int(gap))


# --------------------------------------------------------------------------- #
# Public entry points
# --------------------------------------------------------------------------- #
def _add_ea_slide(prs, data: dict) -> None:
    """Render one EA one-slider onto a new slide of `prs`."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # White background.
    bg = _rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WHITE, rounded=False)
    bg.line.fill.background()

    _header(slide, data)

    # Column geometry.
    margin = Inches(0.3)
    gap = Inches(0.2)
    col_w = (SLIDE_W - 2 * margin - 2 * gap) / 3
    col_w = Emu(int(col_w))
    top = Inches(1.08)
    card_gap = Inches(0.12)
    col_h = Emu(int(SLIDE_H) - int(top) - int(Inches(0.26)))  # usable height

    left_x = margin
    center_x = Emu(int(margin) + int(col_w) + int(gap))
    right_x = Emu(int(margin) + 2 * int(col_w) + 2 * int(gap))

    # --- LEFT column (adaptive: skipped bundle/finite cards free their space) ---
    bundles = data.get("bundles", []) or []
    finite = data.get("finite_licenses", []) or []

    y = top
    h1 = Inches(1.68)
    _contract_details_card(slide, left_x, y, col_w, h1, data)
    y = Emu(int(y) + int(h1) + int(card_gap))
    remaining = Emu(int(col_h) - int(h1) - int(card_gap))

    if bundles and finite:
        # Bundle card sized to its pill count, finite table takes the rest.
        need = Inches(0.5 + 0.42 * min(len(bundles), 4))
        h2 = Emu(max(int(Inches(1.0)), min(int(need), int(Inches(2.2)))))
        _bundle_card(slide, left_x, y, col_w, h2, data)
        y = Emu(int(y) + int(h2) + int(card_gap))
        h3 = Emu(int(remaining) - int(h2) - int(card_gap))
        _finite_card(slide, left_x, y, col_w, h3, data)
    elif finite:
        _finite_card(slide, left_x, y, col_w, remaining, data)
    elif bundles:
        _bundle_card(slide, left_x, y, col_w, remaining, data)
    else:
        ix, iy, iw, ih = _card(slide, left_x, y, col_w, remaining,
                               "Licenses & Bundles")
        _empty_note(slide, ix, iy, iw, ih, "No license or bundle data provided")

    # --- CENTER column ---
    y = top
    hc1 = Inches(3.4)
    _trend_card(slide, center_x, y, col_w, hc1, data)
    y = Emu(int(y) + int(hc1) + int(card_gap))
    hc2 = Inches(2.68)
    _usage_data_card(slide, center_x, y, col_w, hc2, data)

    # --- RIGHT column ---
    y = top
    hr1 = Inches(1.9)
    _locations_card(slide, right_x, y, col_w, hr1, data)
    y = Emu(int(y) + int(hr1) + int(card_gap))
    hr2 = Inches(1.9)
    _version_card(slide, right_x, y, col_w, hr2, data)
    y = Emu(int(y) + int(hr2) + int(card_gap))
    hr3 = Inches(1.18)
    _credits_card(slide, right_x, y, col_w, hr3, data)
    y = Emu(int(y) + int(hr3) + int(card_gap))
    hr4 = Inches(0.86)
    _support_card(slide, right_x, y, col_w, hr4, data)


def _new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def build_slide(data: dict, insights: list[dict] | None = None) -> BytesIO:
    """Build the one-slider (plus an optional insights slide) as a .pptx."""
    prs = _new_presentation()
    _add_ea_slide(prs, data)
    if insights:
        _add_insights_slide(prs, data, insights)
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def build_deck(accounts: list[tuple[dict, list[dict] | None]]) -> BytesIO:
    """Build one deck with a slide (plus optional insights slide) per account.

    `accounts` is a list of (data, insights-or-None) tuples."""
    prs = _new_presentation()
    for data, insights in accounts:
        _add_ea_slide(prs, data)
        if insights:
            _add_insights_slide(prs, data, insights)
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
